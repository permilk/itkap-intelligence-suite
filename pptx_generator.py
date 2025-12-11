"""
ITKAP Intelligence Suite v4.0 - Generador de PowerPoint
Estructura Kenneth:
- 1 portada
- 2 slides generales 
- N slides individuales (uno por colaborador)
- 1 cierre
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import plotly.io as pio
import plotly.graph_objects as go
from datetime import datetime

def generar_presentacion(figuras, df_empleados, df_competencias, stats, config):
    """Genera PPT con estructura específica de Kenneth"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Portada
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_portada(slide, config, len(df_empleados), len(df_competencias.columns))
    
    # General 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_general_1(slide, figuras, df_empleados, df_competencias)
    
    # General 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_general_2(slide, figuras)
    
    # Individuales
    for idx, row in df_empleados.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        crear_individual(slide, row['Nombre'], row['Promedio'], idx, df_empleados, df_competencias)
    
    # Cierre
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    crear_cierre(slide, config)
    
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def crear_portada(slide, config, n_emp, n_comp):
    # Fondo azul
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(14, 27, 46)
    bg.line.fill.background()
    
    # Título
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    box.text_frame.text = "Análisis de Competencias"
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Empresa
    empresa = config.get('nombre_empresa', '')
    if empresa:
        box2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
        box2.text_frame.text = empresa
        p2 = box2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(242, 114, 0)
    
    # Métricas
    box3 = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.6))
    box3.text_frame.text = f"{n_emp} Colaboradores | {n_comp} Competencias"
    p3 = box3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(226, 232, 240)

def crear_general_1(slide, figs, df_emp, df_comp):
    # Título
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    box.text_frame.text = "Vista Organizacional"
    p = box.text_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(14, 27, 46)
    
    # KPIs
    prom = df_emp['Promedio'].mean()
    mejor = df_emp.loc[df_emp['Promedio'].idxmax(), 'Nombre']
    
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.8))
    box2.text_frame.text = f"👥 {len(df_emp)} Colaboradores  |  ⭐ {prom:.2f}  |  🏆 {mejor}"
    p2 = box2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(14)
    
    # Gráficos
    if "04_Competencias_Organizacionales" in figs:
        img = pio.to_image(figs["04_Competencias_Organizacionales"], format='png', width=800, height=500, scale=2)
        slide.shapes.add_picture(BytesIO(img), Inches(0.5), Inches(2), width=Inches(4.5), height=Inches(3))
    
    if "03_Distribucion" in figs:
        img = pio.to_image(figs["03_Distribucion"], format='png', width=800, height=500, scale=2)
        slide.shapes.add_picture(BytesIO(img), Inches(5.2), Inches(2), width=Inches(4.3), height=Inches(3))

def crear_general_2(slide, figs):
    # Título
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    box.text_frame.text = "Mapas de Calor"
    p = box.text_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(14, 27, 46)
    
    # Mapa promedios
    if "01b_Mapa_Calor_Promedios" in figs:
        img = pio.to_image(figs["01b_Mapa_Calor_Promedios"], format='png', width=1600, height=400, scale=2)
        slide.shapes.add_picture(BytesIO(img), Inches(0.5), Inches(1.4), width=Inches(9), height=Inches(1.5))
    
    # Mapa colaboradores  
    if "01_Mapa_de_Calor_Colaboradores" in figs:
        img = pio.to_image(figs["01_Mapa_de_Calor_Colaboradores"], format='png', width=1600, height=900, scale=2)
        slide.shapes.add_picture(BytesIO(img), Inches(0.5), Inches(3.5), width=Inches(9), height=Inches(3.2))

def crear_individual(slide, nombre, prom, idx, df_emp, df_comp):
    # Header
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(14, 27, 46)
    bg.line.fill.background()
    
    # Nombre
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6), Inches(0.6))
    box.text_frame.text = f"Perfil: {nombre}"
    p = box.text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Ranking
    ranking = (df_emp['Promedio'] > prom).sum() + 1
    box2 = slide.shapes.add_textbox(Inches(7), Inches(0.3), Inches(2.5), Inches(0.6))
    box2.text_frame.text = f"#{ranking} | ⭐ {prom:.2f}"
    p2 = box2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(242, 114, 0)
    
    # Mejor/Peor
    comp_emp = df_comp.iloc[idx]
    mejor_c = comp_emp.idxmax()
    mejor_v = comp_emp.max()
    peor_c = comp_emp.idxmin()
    peor_v = comp_emp.min()
    
    box3 = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(0.5))
    box3.text_frame.text = f"🏆 Mejor: {mejor_c} ({mejor_v:.2f})  |  ⚠️ Mejorar: {peor_c} ({peor_v:.2f})"
    p3 = box3.text_frame.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(14)
    
    # Radar
    fig = go.Figure()
    fig.add_trace(go.Scatterpolar(
        r=comp_emp.values.tolist(),
        theta=df_comp.columns.tolist(),
        fill='toself',
        line=dict(color='#e91e63', width=2),
        fillcolor='rgba(233, 30, 99, 0.3)'
    ))
    fig.update_layout(
        polar=dict(radialaxis=dict(visible=True, range=[0, 5])),
        showlegend=False,
        height=400,
        margin=dict(l=80, r=80, t=40, b=40)
    )
    img = pio.to_image(fig, format='png', width=700, height=700, scale=2)
    slide.shapes.add_picture(BytesIO(img), Inches(0.5), Inches(2.1), width=Inches(4.5), height=Inches(4.5))
    
    # Barras
    comp_sorted = comp_emp.sort_values(ascending=True)
    fig2 = go.Figure()
    fig2.add_trace(go.Bar(
        x=comp_sorted.values,
        y=comp_sorted.index,
        orientation='h',
        marker=dict(color=comp_sorted.values, colorscale='RdYlGn', cmin=0, cmax=5),
        text=[f"{v:.2f}" for v in comp_sorted.values],
        textposition='outside'
    ))
    fig2.update_layout(
        title="Competencias",
        xaxis=dict(range=[0, 5.5]),
        height=600,
        margin=dict(l=180, r=50, t=60, b=50)
    )
    img2 = pio.to_image(fig2, format='png', width=900, height=1200, scale=2)
    slide.shapes.add_picture(BytesIO(img2), Inches(5.2), Inches(2.1), width=Inches(4.3), height=Inches(4.5))

def crear_cierre(slide, config):
    # Fondo
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(14, 27, 46)
    bg.line.fill.background()
    
    # Gracias
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    box.text_frame.text = "Gracias"
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # ITKAP
    box2 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    box2.text_frame.text = "ITKAP Consulting"
    p2 = box2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(242, 114, 0)
